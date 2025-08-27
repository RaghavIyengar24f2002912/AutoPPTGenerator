import streamlit as st
import openai
import anthropic
import google.generativeai as genai
import json
from pptx import Presentation
import io

# --- Page Configuration ---
st.set_page_config(
    page_title="Auto-PPT Generator",
    page_icon="📄",
    layout="centered"
)

# --- Functions ---

def generate_slide_structure(text, guidance, api_key, llm_provider):
    """
    Uses the selected LLM to analyze text and return a structured slide plan.
    """
    prompt = f"""
    You are an expert presentation creator. Your task is to analyze the following text and optional guidance
    to create a structured plan for a PowerPoint presentation. The number of slides should be reasonable
    based on the length and content of the text.

    Your output MUST be a valid JSON array of objects. Each object represents a single slide and must have
    two keys: "title" and "content". The "title" key should have a short, engaging string for the slide title.
    The "content" key should have an array of strings, where each string is a concise bullet point for the slide body.

    User Guidance: "{guidance}"

    Text to Analyze:
    ---
    {text}
    ---
    """
    
    raw_response = ""
    try:
        # --- OPENAI LOGIC ---
        if llm_provider == "OpenAI":
            client = openai.OpenAI(api_key=api_key)
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are an expert presentation creator that only outputs valid JSON."},
                    {"role": "user", "content": prompt}
                ],
                response_format={"type": "json_object"}
            )
            raw_response = response.choices[0].message.content

        # --- ANTHROPIC LOGIC ---
        elif llm_provider == "Anthropic":
            client = anthropic.Anthropic(api_key=api_key)
            response = client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=4096,
                system="You are an expert presentation creator. Your only output must be the requested JSON array of slides.",
                messages=[{"role": "user", "content": prompt}]
            )
            raw_response = response.content[0].text

        # --- GEMINI LOGIC ---
        elif llm_provider == "Google Gemini":
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-1.5-flash')
            generation_config = genai.types.GenerationConfig(
                response_mime_type="application/json"
            )
            response = model.generate_content(prompt, generation_config=generation_config)
            raw_response = response.text

        # --- JSON PARSING (common for all) ---
        slide_data = json.loads(raw_response)
        if isinstance(slide_data, dict) and len(slide_data.keys()) == 1:
             return list(slide_data.values())[0]
        return slide_data

    except Exception as e:
        st.error(f"An error occurred with {llm_provider}: {e}")
        return None

def find_layout_index(prs, layout_name):
    """Finds the index of a layout by its name."""
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name.lower() == layout_name.lower():
            return i
    return None

def create_presentation(slide_structure, template_file):
    """
    Generates a .pptx file from a slide structure and a template.
    """
    try:
        prs = Presentation(template_file)
        
        # Heuristically find the best layouts
        title_layout_idx = find_layout_index(prs, 'Title Slide') or 0
        content_layout_idx = find_layout_index(prs, 'Title and Content') or 1

        # Create Title Slide
        if slide_structure:
            title_slide_data = slide_structure[0]
            slide_layout = prs.slide_layouts[title_layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            
            if slide.shapes.title:
                slide.shapes.title.text = title_slide_data.get('title', 'Presentation Title')
            
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = title_slide_data.get('content', [''])[0] or "Generated with AI"

        # Create Content Slides
        for slide_data in slide_structure[1:]:
            slide_layout = prs.slide_layouts[content_layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get('title', '')

            body_shape = None
            for shape in slide.shapes:
                if shape.name.lower().startswith(('content', 'body', 'text')) and shape.has_text_frame:
                    body_shape = shape
                    break
            
            if body_shape:
                tf = body_shape.text_frame
                tf.clear()
                for point in slide_data.get('content', []):
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0

        # Save to a memory buffer
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream

    except Exception as e:
        st.error(f"Error creating presentation: {e}")
        return None


# --- UI Starts Here ---

st.title("Your Text, Your Style – Auto-Generate a Presentation") # [cite: 1]
st.markdown("Create a presentation from bulk text, markdown, or prose using a template of your choice.") # [cite: 2]

st.header("1. Paste Your Content")
input_text = st.text_area("Enter the text to be converted into slides.", height=250) # [cite: 5, 16]

st.header("2. Provide Guidance (Optional)")
guidance = st.text_input("Enter a one-line instruction for tone or structure (e.g., 'turn into an investor pitch deck').") # [cite: 6]

st.header("3. Upload Your Template")
uploaded_file = st.file_uploader("Upload a .pptx or .potx file.", type=['pptx', 'potx']) # [cite: 8, 18]

st.header("4. Configure Your LLM")
llm_provider = st.selectbox("Choose your LLM Provider", ("OpenAI", "Anthropic", "Google Gemini")) # [cite: 7, 24]

api_key = st.text_input(
    f"Enter your {llm_provider} API key. Your key is not stored or logged.", # [cite: 17, 48]
    type="password"
)

st.divider()

if st.button("✨ Generate Presentation", type="primary", use_container_width=True):
    is_ready = True
    if not input_text:
        st.error("Please paste your content to generate slides.")
        is_ready = False
    if not uploaded_file:
        st.error("Please upload a PowerPoint template file.")
        is_ready = False
    if not api_key:
        st.error(f"Please enter your {llm_provider} API key.")
        is_ready = False
    
    if is_ready:
        with st.spinner(f"Asking {llm_provider} to structure your slides... 🧠"):
            slide_structure = generate_slide_structure(input_text, guidance, api_key, llm_provider)

        if slide_structure:
            st.success("Slide structure generated successfully!")
            
            with st.spinner("Applying your template and creating the presentation... 🎨"):
                presentation_file = create_presentation(slide_structure, uploaded_file)
            
            if presentation_file:
                st.success("Your presentation is ready for download!")
                
                st.download_button(
                    label="📥 Download Presentation",
                    data=presentation_file,
                    file_name="generated_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                ) # [cite: 19]
